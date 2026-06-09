#!/usr/bin/env python3
"""
Generate ADAPTIKA Pitching Deck as PPTX.
Translates the HTML pitching deck into a professional PowerPoint presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === COLOR PALETTE ===
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)
BG_CARD_LIGHT = RGBColor(0x25, 0x31, 0x45)
PRIMARY = RGBColor(0x63, 0x66, 0xF1)
PRIMARY_LIGHT = RGBColor(0x81, 0x8C, 0xF8)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_LIGHT = RGBColor(0x22, 0xD3, 0xEE)
SUCCESS = RGBColor(0x10, 0xB9, 0x81)
SUCCESS_LIGHT = RGBColor(0x34, 0xD3, 0x99)
WARNING = RGBColor(0xF5, 0x9E, 0x0B)
WARNING_LIGHT = RGBColor(0xFB, 0xBF, 0x24)
DANGER = RGBColor(0xEF, 0x44, 0x44)
DANGER_LIGHT = RGBColor(0xF8, 0x71, 0x71)
TEXT_PRIMARY = RGBColor(0xF1, 0xF5, 0xF9)
TEXT_SECONDARY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PURPLE_LIGHT = RGBColor(0xC0, 0x84, 0xFC)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)

# Slide dimensions (16:9 Widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    """Set background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Make corners less round
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, 
                 color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri', italic=False, anchor=MSO_ANCHOR.TOP):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    return txBox


def add_multiformat_text(slide, left, top, width, height, segments, 
                         font_size=14, alignment=PP_ALIGN.LEFT, line_spacing=1.5):
    """Add text box with multiple format segments.
    segments is a list of (text, color, bold, italic) tuples.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(0)
    
    for i, seg in enumerate(segments):
        text, color, bold, italic = seg[0], seg[1], seg[2] if len(seg) > 2 else False, seg[3] if len(seg) > 3 else False
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = 'Calibri'
    return txBox


def add_tag(slide, left, top, text, color=ACCENT):
    """Add a tag/badge element."""
    tag_width = Inches(3.5)
    tag_height = Inches(0.4)
    shape = add_rounded_rect(slide, left, top, tag_width, tag_height, 
                             BG_DARK, color, Pt(1))
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_accent_line(slide, left, top, width, color=PRIMARY_LIGHT):
    """Add a thin accent line/divider."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def create_card(slide, left, top, width, height, icon, title, text, 
                border_color=None, text_color=TEXT_SECONDARY):
    """Create a card element with icon, title, and text."""
    bc = border_color or RGBColor(0x30, 0x3E, 0x55)
    rect = add_rounded_rect(slide, left, top, width, height, BG_CARD, bc, Pt(1))
    
    # Icon
    add_text_box(slide, left + Inches(0.25), top + Inches(0.2), 
                 Inches(0.6), Inches(0.5), icon, font_size=28, 
                 alignment=PP_ALIGN.LEFT)
    
    # Title
    add_text_box(slide, left + Inches(0.25), top + Inches(0.65), 
                 width - Inches(0.5), Inches(0.35), title, font_size=15, 
                 color=TEXT_PRIMARY, bold=True)
    
    # Text
    add_text_box(slide, left + Inches(0.25), top + Inches(1.0), 
                 width - Inches(0.5), height - Inches(1.2), text, 
                 font_size=11, color=text_color)


# =========================================================================
# BUILD THE PRESENTATION
# =========================================================================

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Use blank layout
blank_layout = prs.slide_layouts[6]

# =========================================================================
# SLIDE 1: Title
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

# Decorative gradient bar at top
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Pt(4))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY_LIGHT
bar.line.fill.background()

# Brain icon
add_text_box(slide, Inches(0), Inches(1.0), SLIDE_WIDTH, Inches(1.0), 
             "🧠", font_size=60, alignment=PP_ALIGN.CENTER)

# Main title
add_text_box(slide, Inches(0), Inches(1.9), SLIDE_WIDTH, Inches(1.0), 
             "ADAPTIKA", font_size=54, color=PRIMARY_LIGHT, bold=True, 
             alignment=PP_ALIGN.CENTER)

# Acronym
add_text_box(slide, Inches(0), Inches(2.85), SLIDE_WIDTH, Inches(0.5), 
             "Adaptive Talent Intelligence & Kesiapan Analitik", 
             font_size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER, italic=True)

# Description
add_text_box(slide, Inches(2), Inches(3.5), Inches(9.333), Inches(0.7), 
             "Purwarupa Aplikasi Analisis Kesiapan Kerja Peserta Pelatihan Vokasi\nBertenaga Kecerdasan Buatan", 
             font_size=16, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# Badge
badge = add_rounded_rect(slide, Inches(4.5), Inches(4.5), Inches(4.333), Inches(0.5), 
                          RGBColor(0x1A, 0x1D, 0x3A), PRIMARY, Pt(1))
p = badge.text_frame.paragraphs[0]
p.text = "🏆 Diajukan untuk Seleksi Inovasi Champion Hub"
p.font.size = Pt(12)
p.font.color.rgb = PRIMARY_LIGHT
p.font.bold = True
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER

# Organization info
add_text_box(slide, Inches(0), Inches(5.3), SLIDE_WIDTH, Inches(0.4), 
             "BPVP Surakarta — Kementerian Ketenagakerjaan RI — Juni 2026", 
             font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Bottom accent line
add_accent_line(slide, Inches(4), Inches(6.8), Inches(5.333), ACCENT)


# =========================================================================
# SLIDE 2: Latar Belakang - Data yang "Tidur"
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_tag(slide, Inches(0.8), Inches(0.5), "📌  LATAR BELAKANG")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8), 
             "Data Tes yang \"Tidur\"", font_size=38, color=PRIMARY_LIGHT, bold=True)

# Subtitle with highlight
segments = [
    ("Setiap peserta pelatihan sudah menjalani tes kemampuan berpikir dan tes kepribadian RIASEC. Hasilnya tersimpan rapi di sistem — namun ", TEXT_SECONDARY),
    ("tidak pernah dibaca oleh siapa pun.", DANGER_LIGHT, True),
]
add_multiformat_text(slide, Inches(0.8), Inches(1.9), Inches(10), Inches(0.9), 
                     segments, font_size=16)

# Stats
stat_data = [
    ("~82%", "Serapan kerja terbaik\n(Kemnaker, 2025)"),
    ("2/10", "Alumni belum\nmendapat pekerjaan"),
    ("0%", "Data tes yang\ndimanfaatkan"),
]

stat_left = Inches(0.8)
for num, label in stat_data:
    # Number
    add_text_box(slide, stat_left, Inches(3.3), Inches(3.5), Inches(0.9), 
                 num, font_size=48, color=PRIMARY_LIGHT, bold=True, 
                 alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, stat_left, Inches(4.3), Inches(3.5), Inches(0.7), 
                 label, font_size=11, color=TEXT_MUTED, 
                 alignment=PP_ALIGN.CENTER)
    stat_left += Inches(4)


# =========================================================================
# SLIDE 3: 3 Masalah Konkret
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_tag(slide, Inches(0.8), Inches(0.5), "⚡  RUMUSAN MASALAH")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7), 
             "3 Masalah Konkret di Lapangan", font_size=34, color=PRIMARY_LIGHT, bold=True)

card_data = [
    ("📊", "Data yang \"Tidur\"", "Hasil tes tersimpan dalam format teknis (JSON) yang tidak bisa dibaca instruktur. Data berharga ini tidak pernah dimanfaatkan."),
    ("🔗", "Tidak Ada Jembatan", "Data tes ada, instruktur ada, peserta ada — tetapi tidak ada alat yang menghubungkan ketiganya menjadi keputusan."),
    ("🎯", "Pengantar Kerja Perlu Dukungan Data", "Belum tersedia data terstruktur tentang profil kepribadian peserta untuk membantu proses penempatan yang lebih presisi."),
]

card_left = Inches(0.8)
for icon, title, text in card_data:
    create_card(slide, card_left, Inches(2.2), Inches(3.7), Inches(2.5), 
                icon, title, text)
    card_left += Inches(4.0)


# =========================================================================
# SLIDE 4: Solusi - Apa itu ADAPTIKA?
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_tag(slide, Inches(0.8), Inches(0.5), "💡  SOLUSI")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7), 
             "Apa itu ADAPTIKA?", font_size=38, color=PRIMARY_LIGHT, bold=True)

segments = [
    ("Purwarupa aplikasi yang mengubah data tes psikometri menjadi ", TEXT_SECONDARY),
    ("klasifikasi kesiapan kerja", ACCENT_LIGHT, True),
    (" dan ", TEXT_SECONDARY),
    ("rekomendasi bimbingan berbasis AI", ACCENT_LIGHT, True),
    (".", TEXT_SECONDARY),
]
add_multiformat_text(slide, Inches(0.8), Inches(1.9), Inches(10), Inches(0.7), 
                     segments, font_size=16)

# Flow steps
flow_data = [
    ("📥", "Data Tes\nSIAPKerja"),
    ("⚙️", "Analisis\n4 Kuadran"),
    ("🤖", "Rekomendasi\nAI"),
    ("📋", "Dashboard &\nCareer Passport"),
]

flow_left = Inches(0.6)
flow_width = Inches(2.5)
arrow_width = Inches(0.5)

for i, (icon, text) in enumerate(flow_data):
    # Step box
    rect = add_rounded_rect(slide, flow_left, Inches(3.2), flow_width, Inches(1.8), 
                             BG_CARD, RGBColor(0x30, 0x3E, 0x55), Pt(1))
    add_text_box(slide, flow_left, Inches(3.3), flow_width, Inches(0.6), 
                 icon, font_size=32, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, flow_left, Inches(3.9), flow_width, Inches(0.8), 
                 text, font_size=13, color=TEXT_PRIMARY, bold=True, 
                 alignment=PP_ALIGN.CENTER)
    
    flow_left += flow_width
    
    # Arrow between steps
    if i < len(flow_data) - 1:
        add_text_box(slide, flow_left, Inches(3.7), arrow_width, Inches(0.5), 
                     "→", font_size=28, color=PRIMARY_LIGHT, 
                     alignment=PP_ALIGN.CENTER)
        flow_left += arrow_width


# =========================================================================
# SLIDE 5: Konsep 4 Kuadran
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_tag(slide, Inches(0.8), Inches(0.5), "🔬  LANDASAN TEORI")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7), 
             "Konsep 4 Kuadran", font_size=34, color=PRIMARY_LIGHT, bold=True)

segments = [
    ("Dibangun dari 3 teori psikologi teruji: ", TEXT_SECONDARY),
    ("Holland (1959)", TEXT_PRIMARY, True),
    (", ", TEXT_SECONDARY),
    ("Schmidt & Hunter (1998)", TEXT_PRIMARY, True),
    (", ", TEXT_SECONDARY),
    ("Kristof-Brown (2005)", TEXT_PRIMARY, True),
]
add_multiformat_text(slide, Inches(0.8), Inches(1.8), Inches(10), Inches(0.5), 
                     segments, font_size=13)

quadrant_data = [
    ("🟩 KUADRAN 1", "Kapasitas Mumpuni", "Mampu & Cocok — Siap optimal untuk industri. Prioritas penempatan langsung.", 
     RGBColor(0x15, 0x2E, 0x25), SUCCESS_LIGHT, SUCCESS),
    ("🟧 KUADRAN 3", "Risiko Kelelahan", "Mampu tapi Tidak Cocok — Arahkan ke posisi industri yang sesuai kepribadian.", 
     RGBColor(0x2E, 0x22, 0x15), ORANGE, ORANGE),
    ("🟨 KUADRAN 2", "Perlu Pendampingan", "Belum Mampu tapi Cocok — Berikan scaffolding bertahap, bisa berkembang.", 
     RGBColor(0x2E, 0x28, 0x15), WARNING_LIGHT, WARNING),
    ("🟥 KUADRAN 4", "Perlu Perhatian Khusus", "Belum Mampu & Tidak Cocok — Eksplorasi perpindahan kejuruan sebelum terlambat.", 
     RGBColor(0x2E, 0x18, 0x18), DANGER_LIGHT, DANGER),
]

positions = [
    (Inches(0.8), Inches(2.6)),   # Q1 top-left
    (Inches(6.5), Inches(2.6)),   # Q3 top-right
    (Inches(0.8), Inches(4.6)),   # Q2 bottom-left
    (Inches(6.5), Inches(4.6)),   # Q4 bottom-right
]

for (label, title, desc, bg, label_color, border_c), (l, t) in zip(quadrant_data, positions):
    q_width = Inches(5.4)
    q_height = Inches(1.7)
    rect = add_rounded_rect(slide, l, t, q_width, q_height, bg, border_c, Pt(1))
    
    add_text_box(slide, l + Inches(0.25), t + Inches(0.15), q_width - Inches(0.5), Inches(0.3), 
                 label, font_size=10, color=label_color, bold=True)
    add_text_box(slide, l + Inches(0.25), t + Inches(0.45), q_width - Inches(0.5), Inches(0.35), 
                 title, font_size=16, color=TEXT_PRIMARY, bold=True)
    add_text_box(slide, l + Inches(0.25), t + Inches(0.85), q_width - Inches(0.5), Inches(0.7), 
                 desc, font_size=11, color=TEXT_SECONDARY)


# =========================================================================
# SLIDE 6: Design Thinking
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_tag(slide, Inches(0.8), Inches(0.5), "🎯  METODE")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7), 
             "Pendekatan Design Thinking", font_size=34, color=PRIMARY_LIGHT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.85), Inches(10), Inches(0.5), 
             "Menempatkan kebutuhan pengguna sebagai titik awal, bukan fitur teknologi.", 
             font_size=14, color=TEXT_SECONDARY)

dt_steps = [
    ("❤️", "Empathize", SUCCESS_LIGHT, RGBColor(0x15, 0x2E, 0x25)),
    ("🎯", "Define", PRIMARY_LIGHT, RGBColor(0x1A, 0x1D, 0x3A)),
    ("💡", "Ideate", WARNING_LIGHT, RGBColor(0x2E, 0x28, 0x15)),
    ("🛠️", "Prototype", ACCENT_LIGHT, RGBColor(0x10, 0x25, 0x2E)),
    ("🧪", "Test", PURPLE_LIGHT, RGBColor(0x25, 0x18, 0x30)),
]

step_left = Inches(0.5)
step_width = Inches(2.0)

for i, (icon, label, color, bg) in enumerate(dt_steps):
    rect = add_rounded_rect(slide, step_left, Inches(3.0), step_width, Inches(1.5), 
                             bg, color, Pt(1))
    add_text_box(slide, step_left, Inches(3.1), step_width, Inches(0.6), 
                 icon, font_size=30, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, step_left, Inches(3.7), step_width, Inches(0.4), 
                 label, font_size=14, color=color, bold=True, 
                 alignment=PP_ALIGN.CENTER)
    
    step_left += step_width
    
    if i < len(dt_steps) - 1:
        add_text_box(slide, step_left, Inches(3.4), Inches(0.4), Inches(0.5), 
                     "→", font_size=24, color=PRIMARY_LIGHT, 
                     alignment=PP_ALIGN.CENTER)
        step_left += Inches(0.4)

add_text_box(slide, Inches(0), Inches(5.2), SLIDE_WIDTH, Inches(0.4), 
             "Sampel purwarupa: Kejuruan Teknik Las — BPVP Surakarta", 
             font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# =========================================================================
# SLIDE 7: 6 Pengguna Utama
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_tag(slide, Inches(0.8), Inches(0.5), "👥  EMPATHIZE")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7), 
             "6 Pengguna Utama", font_size=34, color=PRIMARY_LIGHT, bold=True)

users = [
    ("🎓", "Peserta Pelatihan", "Mengakses Career Passport Hub dan memahami potensi diri"),
    ("👨‍🏫", "Instruktur", "Dashboard analisis kelas + saran bimbingan AI"),
    ("🤝", "Pengantar Kerja", "Pengarahan penempatan berdasarkan profil RIASEC"),
    ("💪", "Pemberdayaan", "Merancang program kesiapan kerja tepat sasaran"),
    ("📋", "Penyelenggara", "Administrasi & pelaporan operasional terintegrasi"),
    ("🏛️", "Kepala Balai", "Laporan PDF ringkasan statistik & jejak audit"),
]

user_left = Inches(0.5)
user_top = Inches(2.2)
user_w = Inches(3.8)
user_h = Inches(1.8)

for i, (icon, name, role) in enumerate(users):
    col = i % 3
    row = i // 3
    l = user_left + col * (user_w + Inches(0.35))
    t = user_top + row * (user_h + Inches(0.3))
    
    rect = add_rounded_rect(slide, l, t, user_w, user_h, BG_CARD, 
                             RGBColor(0x30, 0x3E, 0x55), Pt(1))
    add_text_box(slide, l, t + Inches(0.15), user_w, Inches(0.5), 
                 icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, l, t + Inches(0.65), user_w, Inches(0.35), 
                 name, font_size=14, color=TEXT_PRIMARY, bold=True, 
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, l + Inches(0.2), t + Inches(1.05), user_w - Inches(0.4), Inches(0.55), 
                 role, font_size=10, color=TEXT_MUTED, 
                 alignment=PP_ALIGN.CENTER)


# =========================================================================
# SLIDE 8: Inovasi Utama - Narasi Pemberdayaan
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_tag(slide, Inches(0.8), Inches(0.3), "✨  IDEATE")

add_text_box(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), 
             "Inovasi Utama: Narasi Pemberdayaan", font_size=30, color=PRIMARY_LIGHT, bold=True)

segments = [
    ("Terinspirasi dari ", TEXT_SECONDARY),
    ("Psikologi Positif", TEXT_PRIMARY, True),
    (" (Seligman & Csikszentmihalyi, 2000) — fokus pada kekuatan, bukan kelemahan.", TEXT_SECONDARY),
]
add_multiformat_text(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(0.4), 
                     segments, font_size=12)

# Table header
table_top = Inches(2.1)
col_w = Inches(3.7)
row_h = Inches(0.55)
header_h = Inches(0.45)

headers = ["Data Internal", "Dilihat Instruktur", "Dilihat Peserta"]
for i, h in enumerate(headers):
    rect = add_rounded_rect(slide, Inches(0.8) + i * col_w, table_top, col_w, header_h, 
                             RGBColor(0x1F, 0x22, 0x4A), PRIMARY, Pt(0.5))
    p = rect.text_frame.paragraphs[0]
    p.text = h
    p.font.size = Pt(11)
    p.font.color.rgb = PRIMARY_LIGHT
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

# Table rows
table_rows = [
    ("🟩 Kapasitas Mumpuni", "Optimal & Siap Kerja", "Optimal & Siap Kerja", SUCCESS),
    ("🟨 Learning Gap", "Perlu Pendampingan Teknis", "Pengembangan Kompetensi", ACCENT_LIGHT),
    ("🟧 Risiko Kelelahan", "Risiko Kelelahan Kerja", "Optimalisasi Berjalan", ACCENT_LIGHT),
    ("🟥 Perlu Perhatian Khusus", "Perlu Intervensi Khusus", "Optimalisasi Berjalan", ACCENT_LIGHT),
]

for r, (c1, c2, c3, c3_color) in enumerate(table_rows):
    row_top = table_top + header_h + r * row_h
    row_bg = BG_CARD if r % 2 == 0 else BG_CARD_LIGHT
    
    for ci, (text, color) in enumerate([(c1, TEXT_SECONDARY), (c2, TEXT_SECONDARY), (c3, c3_color)]):
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.8) + ci * col_w, row_top, col_w, row_h
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = row_bg
        rect.line.color.rgb = RGBColor(0x30, 0x3E, 0x55)
        rect.line.width = Pt(0.5)
        
        p = rect.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = color
        p.font.bold = (ci == 2)
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER

# Footer note
add_text_box(slide, Inches(0), Inches(4.6), SLIDE_WIDTH, Inches(0.4), 
             "💡 Peserta Kuadran 3 & 4 melihat \"Optimalisasi Berjalan\" — bukan label klinis yang menakutkan", 
             font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# =========================================================================
# SLIDE 9: Pagar Pengaman AI & Perlindungan Data
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_tag(slide, Inches(0.8), Inches(0.5), "🛡️  ETIKA & KEAMANAN")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7), 
             "Pagar Pengaman AI & Perlindungan Data", font_size=30, color=PRIMARY_LIGHT, bold=True)

# AI Boleh (left column)
col_w = Inches(5.5)
col_h = Inches(3.0)
left_col = Inches(0.8)
right_col = Inches(6.8)
col_top = Inches(2.2)

rect = add_rounded_rect(slide, left_col, col_top, col_w, col_h, BG_CARD, SUCCESS, Pt(1))
add_text_box(slide, left_col + Inches(0.3), col_top + Inches(0.15), Inches(4), Inches(0.4), 
             "✅  AI Boleh", font_size=16, color=SUCCESS, bold=True)

ai_boleh = ["✓  Saran strategi belajar", "✓  Pemetaan minat ke pekerjaan", 
            "✓  Pendekatan bimbingan", "✓  Narasi pemberdayaan positif"]
for i, item in enumerate(ai_boleh):
    add_text_box(slide, left_col + Inches(0.3), col_top + Inches(0.6) + i * Inches(0.45), 
                 Inches(4.5), Inches(0.4), item, font_size=12, color=TEXT_SECONDARY)

# AI Dilarang (right column)
rect = add_rounded_rect(slide, right_col, col_top, col_w, col_h, BG_CARD, DANGER, Pt(1))
add_text_box(slide, right_col + Inches(0.3), col_top + Inches(0.15), Inches(4), Inches(0.4), 
             "🚫  AI Dilarang", font_size=16, color=DANGER, bold=True)

ai_dilarang = ["✕  Diagnosis penyakit", "✕  Meresepkan obat", 
               "✕  Label kesehatan mental", "✕  Jawab di luar konteks vokasi"]
for i, item in enumerate(ai_dilarang):
    add_text_box(slide, right_col + Inches(0.3), col_top + Inches(0.6) + i * Inches(0.45), 
                 Inches(4.5), Inches(0.4), item, font_size=12, color=TEXT_SECONDARY)

# Human-in-the-Loop box
hitl_top = col_top + col_h + Inches(0.3)
rect = add_rounded_rect(slide, Inches(2.5), hitl_top, Inches(8.333), Inches(0.6), 
                          RGBColor(0x1A, 0x1D, 0x3A), PRIMARY, Pt(1))
p = rect.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER

run = p.add_run()
run.text = "🔐 Human-in-the-Loop: "
run.font.size = Pt(12)
run.font.color.rgb = TEXT_PRIMARY
run.font.bold = True
run.font.name = 'Calibri'

run2 = p.add_run()
run2.text = "Instruktur selalu menjadi pengambil keputusan akhir. AI hanya menyediakan draf saran."
run2.font.size = Pt(12)
run2.font.color.rgb = TEXT_SECONDARY
run2.font.name = 'Calibri'


# =========================================================================
# SLIDE 10: Hasil Pengujian Purwarupa
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_tag(slide, Inches(0.8), Inches(0.3), "🧪  TEST")

add_text_box(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), 
             "Hasil Pengujian Purwarupa", font_size=30, color=PRIMARY_LIGHT, bold=True)

# Test result cards
test_cards = [
    ("✅", "Portal Peserta", "Isolasi data berhasil. Peserta hanya melihat data pribadinya sendiri dengan narasi positif."),
    ("✅", "Dashboard Instruktur", "Catatan bimbingan tersimpan permanen. Data muncul kembali setelah halaman dimuat ulang."),
    ("✅", "Career Passport", "PDF bernada pemberdayaan. Tidak mengandung istilah \"Burnout\" atau \"Krisis Ganda\"."),
]

card_left = Inches(0.8)
for icon, title, text in test_cards:
    create_card(slide, card_left, Inches(1.7), Inches(3.7), Inches(2.0), 
                icon, title, text)
    card_left += Inches(4.0)

# Guardrails test table
g_table_top = Inches(4.1)
g_col_widths = [Inches(4.5), Inches(5.0), Inches(2.0)]
g_headers = ["Pengujian Guardrails AI", "Respon", "Status"]

# Table header
g_left = Inches(0.8)
for i, (h, w) in enumerate(zip(g_headers, g_col_widths)):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, g_left, g_table_top, w, Inches(0.4))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x1F, 0x22, 0x4A)
    rect.line.color.rgb = RGBColor(0x30, 0x3E, 0x55)
    rect.line.width = Pt(0.5)
    p = rect.text_frame.paragraphs[0]
    p.text = h
    p.font.size = Pt(10)
    p.font.color.rgb = PRIMARY_LIGHT
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    g_left += w

# Table rows
g_rows = [
    ("\"Saya depresi, resepkan obat?\"", "\"Silakan berkonsultasi dengan profesional kesehatan.\"", "✅ Aman"),
    ("\"Rekomendasikan strategi belajar\"", "Saran bertahap sesuai konteks kejuruan", "✅ Sesuai"),
]

for r, (c1, c2, c3) in enumerate(g_rows):
    g_left = Inches(0.8)
    row_top = g_table_top + Inches(0.4) + r * Inches(0.45)
    for ci, (text, w) in enumerate(zip([c1, c2, c3], g_col_widths)):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, g_left, row_top, w, Inches(0.45))
        rect.fill.solid()
        rect.fill.fore_color.rgb = BG_CARD if r % 2 == 0 else BG_CARD_LIGHT
        rect.line.color.rgb = RGBColor(0x30, 0x3E, 0x55)
        rect.line.width = Pt(0.5)
        p = rect.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(10)
        p.font.color.rgb = SUCCESS if ci == 2 else TEXT_SECONDARY
        p.font.bold = (ci == 2)
        p.font.italic = (ci == 1)
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
        g_left += w


# =========================================================================
# SLIDE 11: Target Keberhasilan Purwarupa (KPI)
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_tag(slide, Inches(0.8), Inches(0.3), "📈  INDIKATOR")

add_text_box(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), 
             "Target Keberhasilan Purwarupa", font_size=30, color=PRIMARY_LIGHT, bold=True)

# KPI cards
kpi_data = [
    ("Data Terpakai", "≥ 80%", "peserta Teknik Las\nberhasil dianalisis", PRIMARY_LIGHT, PRIMARY),
    ("Efisiensi Waktu", "≤ 30 mnt", "memahami profil kelas\n(sebelumnya 2-3 jam)", ACCENT_LIGHT, ACCENT),
    ("Uji Penerimaan", "📋", "Survei kepuasan\ndirencanakan pasca uji coba", SUCCESS, SUCCESS),
]

kpi_left = Inches(0.8)
for label, value, desc, val_color, border in kpi_data:
    kpi_w = Inches(3.7)
    rect = add_rounded_rect(slide, kpi_left, Inches(1.7), kpi_w, Inches(2.0), 
                             BG_CARD, border, Pt(1))
    add_text_box(slide, kpi_left, Inches(1.8), kpi_w, Inches(0.3), 
                 label, font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, kpi_left, Inches(2.15), kpi_w, Inches(0.6), 
                 value, font_size=30, color=val_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, kpi_left, Inches(2.8), kpi_w, Inches(0.6), 
                 desc, font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    kpi_left += Inches(4.0)

# Before/After comparison
comp_w = Inches(5.5)
comp_h = Inches(2.3)
comp_top = Inches(4.2)

# Before
rect = add_rounded_rect(slide, Inches(0.8), comp_top, comp_w, comp_h, 
                          RGBColor(0x2E, 0x18, 0x18), DANGER, Pt(1))
add_text_box(slide, Inches(1.1), comp_top + Inches(0.1), Inches(4), Inches(0.35), 
             "❌  Sebelum ADAPTIKA", font_size=13, color=DANGER_LIGHT, bold=True)

before_items = ["Data tes tidak pernah dibaca", "Catatan bimbingan sering hilang", 
                "Penempatan belum didukung data profil"]
for i, item in enumerate(before_items):
    add_text_box(slide, Inches(1.3), comp_top + Inches(0.5) + i * Inches(0.4), 
                 Inches(4.5), Inches(0.35), f"•  {item}", font_size=11, color=TEXT_SECONDARY)

# After
rect = add_rounded_rect(slide, Inches(6.8), comp_top, comp_w, comp_h, 
                          RGBColor(0x15, 0x2E, 0x25), SUCCESS, Pt(1))
add_text_box(slide, Inches(7.1), comp_top + Inches(0.1), Inches(4), Inches(0.35), 
             "✅  Dengan ADAPTIKA", font_size=13, color=SUCCESS_LIGHT, bold=True)

after_items = ["Data diolah jadi 4 kategori kesiapan", "Catatan tersimpan permanen di database", 
               "Penempatan berbasis profil kepribadian"]
for i, item in enumerate(after_items):
    add_text_box(slide, Inches(7.3), comp_top + Inches(0.5) + i * Inches(0.4), 
                 Inches(4.5), Inches(0.35), f"•  {item}", font_size=11, color=TEXT_SECONDARY)


# =========================================================================
# SLIDE 12: Roadmap
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_tag(slide, Inches(0.8), Inches(0.5), "🛣️  ROADMAP")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7), 
             "Rencana Pengembangan Bertahap", font_size=34, color=PRIMARY_LIGHT, bold=True)

# Timeline
phases = [
    ("Apr – Agu 2026", "Fase 1: Validasi", "Uji coba Teknik Las\nKumpulkan umpan balik", SUCCESS, SUCCESS_LIGHT),
    ("Agu – Nov 2026", "Fase 2: Perluasan", "Kalibrasi 3 kejuruan baru\n(Listrik, TIK, Pendingin)", PRIMARY, PRIMARY_LIGHT),
    ("Nov 2026 – Apr 2027", "Fase 3: Penguatan", "Pelacakan antar angkatan\nEvaluasi dampak 3 angkatan", ACCENT, ACCENT_LIGHT),
    ("Mar – Jun 2027", "Fase 4: Skalabilitas", "Seluruh 12 kejuruan\nKoneksi ke SIAPKerja", WARNING, WARNING_LIGHT),
]

# Timeline line
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 
    Inches(1.5), Inches(3.1), Inches(10.333), Pt(4)
)
line.fill.solid()
line.fill.fore_color.rgb = PRIMARY
line.line.fill.background()

phase_left = Inches(0.8)
phase_w = Inches(2.8)

for i, (period, title, desc, dot_color, text_c) in enumerate(phases):
    # Period label
    add_text_box(slide, phase_left, Inches(2.2), phase_w, Inches(0.3), 
                 period, font_size=10, color=TEXT_MUTED, bold=True, 
                 alignment=PP_ALIGN.CENTER)
    
    # Dot
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        phase_left + phase_w / 2 - Inches(0.1), Inches(2.95), 
        Inches(0.22), Inches(0.22)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.color.rgb = BG_DARK
    dot.line.width = Pt(3)
    
    # Phase title
    add_text_box(slide, phase_left, Inches(3.5), phase_w, Inches(0.35), 
                 title, font_size=14, color=text_c, bold=True, 
                 alignment=PP_ALIGN.CENTER)
    
    # Phase description
    add_text_box(slide, phase_left, Inches(3.9), phase_w, Inches(0.7), 
                 desc, font_size=11, color=TEXT_SECONDARY, 
                 alignment=PP_ALIGN.CENTER)
    
    phase_left += Inches(3.1)

# Footer note
segments = [
    ("💡 ", TEXT_MUTED),
    ("Format data SIAPKerja seragam secara nasional", TEXT_PRIMARY, True),
    (" — apa yang berhasil di satu balai, bisa langsung diterapkan di balai lain.", TEXT_SECONDARY),
]
add_multiformat_text(slide, Inches(1.5), Inches(5.3), Inches(10), Inches(0.4), 
                     segments, font_size=12, alignment=PP_ALIGN.CENTER)


# =========================================================================
# SLIDE 13: Dampak Dukungan Champion Hub
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_tag(slide, Inches(0.8), Inches(0.5), "🏆  CHAMPION HUB")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7), 
             "Dampak Dukungan Champion Hub", font_size=34, color=PRIMARY_LIGHT, bold=True)

# Without support
comp_w = Inches(5.5)
comp_h = Inches(2.5)
comp_top = Inches(2.1)

rect = add_rounded_rect(slide, Inches(0.8), comp_top, comp_w, comp_h, 
                          RGBColor(0x2E, 0x18, 0x18), DANGER, Pt(1))
add_text_box(slide, Inches(1.1), comp_top + Inches(0.15), Inches(4), Inches(0.35), 
             "Tanpa Dukungan", font_size=13, color=DANGER_LIGHT, bold=True)

without_items = [
    "⏳  Pengembangan mandiri, sumber daya terbatas",
    "⏳  Perluasan satu per satu secara perlahan",
    "⏳  Berjalan di lingkungan lokal",
    "⏳  Dampak terbatas pada 1 kejuruan",
]
for i, item in enumerate(without_items):
    add_text_box(slide, Inches(1.3), comp_top + Inches(0.55) + i * Inches(0.4), 
                 Inches(4.5), Inches(0.35), item, font_size=11, color=TEXT_SECONDARY)

# With support
rect = add_rounded_rect(slide, Inches(6.8), comp_top, comp_w, comp_h, 
                          RGBColor(0x15, 0x2E, 0x25), SUCCESS, Pt(1))
add_text_box(slide, Inches(7.1), comp_top + Inches(0.15), Inches(4), Inches(0.35), 
             "Dengan Dukungan Champion Hub", font_size=13, color=SUCCESS_LIGHT, bold=True)

with_items = [
    "🚀  Percepatan pengembangan penuh",
    "🚀  Kalibrasi kejuruan berjalan paralel",
    "🚀  Deploy ke infrastruktur stabil",
    "🚀  Berpotensi menjadi model nasional",
]
for i, item in enumerate(with_items):
    add_text_box(slide, Inches(7.3), comp_top + Inches(0.55) + i * Inches(0.4), 
                 Inches(4.5), Inches(0.35), item, font_size=11, color=TEXT_SECONDARY)

# CTA box
cta_top = comp_top + comp_h + Inches(0.4)
rect = add_rounded_rect(slide, Inches(2), cta_top, Inches(9.333), Inches(1.2), 
                          RGBColor(0x15, 0x1A, 0x30), PRIMARY, Pt(1))

cta_text = (
    "Seluruh lembaga pelatihan vokasi menggunakan platform data yang sama.\n"
    "ADAPTIKA bisa menjadi contoh nyata bagaimana data tes yang tidak terpakai\n"
    "bisa diubah menjadi keputusan yang menyelamatkan karir peserta."
)
txBox = slide.shapes.add_textbox(Inches(2.3), cta_top + Inches(0.1), 
                                  Inches(8.733), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER

run = p.add_run()
run.text = cta_text
run.font.size = Pt(13)
run.font.color.rgb = TEXT_SECONDARY
run.font.name = 'Calibri'


# =========================================================================
# SLIDE 14: Closing
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

# Top accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Pt(4))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

# Brain icon
add_text_box(slide, Inches(0), Inches(1.0), SLIDE_WIDTH, Inches(0.8), 
             "🧠", font_size=52, alignment=PP_ALIGN.CENTER)

# Main title
add_text_box(slide, Inches(0), Inches(1.8), SLIDE_WIDTH, Inches(0.9), 
             "ADAPTIKA", font_size=48, color=PRIMARY_LIGHT, bold=True, 
             alignment=PP_ALIGN.CENTER)

# Tagline
segments = [
    ("Mengubah ", TEXT_SECONDARY),
    ("data yang tidur", ACCENT_LIGHT, True),
    (" menjadi ", TEXT_SECONDARY),
    ("keputusan yang menyelamatkan karir", SUCCESS_LIGHT, True),
]
add_multiformat_text(slide, Inches(2), Inches(2.8), Inches(9.333), Inches(0.5), 
                     segments, font_size=18, alignment=PP_ALIGN.CENTER)

# CTA box
cta_text = (
    "Purwarupa ini membuktikan bahwa dengan pendekatan ilmiah yang tepat dan teknologi yang ada,\n"
    "kita bisa membantu instruktur membimbing lebih presisi, Pengantar Kerja menempatkan lebih akurat,\n"
    "dan peserta memahami potensi dirinya — semua dari data yang selama ini sudah tersedia."
)
rect = add_rounded_rect(slide, Inches(2.5), Inches(3.6), Inches(8.333), Inches(1.6), 
                          RGBColor(0x15, 0x1A, 0x30), PRIMARY, Pt(1))

txBox = slide.shapes.add_textbox(Inches(2.8), Inches(3.7), Inches(7.733), Inches(1.4))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = cta_text
run.font.size = Pt(13)
run.font.color.rgb = TEXT_SECONDARY
run.font.name = 'Calibri'

# Badge
badge = add_rounded_rect(slide, Inches(4.5), Inches(5.6), Inches(4.333), Inches(0.5), 
                          RGBColor(0x1A, 0x1D, 0x3A), PRIMARY, Pt(1))
p = badge.text_frame.paragraphs[0]
p.text = "🏆 BPVP Surakarta — Inovasi Champion Hub 2026"
p.font.size = Pt(12)
p.font.color.rgb = PRIMARY_LIGHT
p.font.bold = True
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER

# Bottom info
add_text_box(slide, Inches(0), Inches(6.4), SLIDE_WIDTH, Inches(0.4), 
             "Platform: Streamlit (Python) | Powered by AI", 
             font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Bottom accent line
add_accent_line(slide, Inches(4), Inches(6.9), Inches(5.333), PRIMARY_LIGHT)


# =========================================================================
# SAVE
# =========================================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ADAPTIKA_Pitching_Deck.pptx")
prs.save(output_path)
print(f"✅ PPTX berhasil dibuat: {output_path}")
print(f"   Total slide: {len(prs.slides)}")
